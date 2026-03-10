"""
Attachment text extraction for email indexing.

Extracts searchable text from common attachment types:
- PDF (.pdf)
- Word (.docx)
- Excel (.xlsx)
- PowerPoint (.pptx)
- Plain text (.txt, .csv, .log, .json, .xml)
- HTML (.html, .htm)

Size guard: skips attachments >25MB.
No OCR — only embedded text is extracted.
"""

import io
import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

MAX_ATTACHMENT_SIZE = 25 * 1024 * 1024  # 25MB

# Map of MIME types to extractor functions
_MIME_EXTRACTORS = {}

# Map of file extensions to extractor functions (fallback for octet-stream)
_EXT_EXTRACTORS = {}

# Supported MIME types (for can_extract checks)
SUPPORTED_MIME_TYPES = {
    'application/pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'application/msword',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    'application/vnd.ms-excel',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    'application/vnd.ms-powerpoint',
    'text/plain', 'text/csv', 'text/html',
}

SUPPORTED_EXTENSIONS = {
    '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt',
    '.txt', '.csv', '.log', '.json', '.xml', '.html', '.htm', '.md',
}


def can_extract(content_type: str, filename: str) -> bool:
    """Check if we can extract text from this file type."""
    if content_type and content_type.lower() in SUPPORTED_MIME_TYPES:
        return True
    if filename:
        ext = Path(filename).suffix.lower()
        return ext in SUPPORTED_EXTENSIONS
    return False


def extract_text(filename: str, content: bytes, content_type: str) -> str:
    """Extract searchable text from attachment bytes.

    Returns extracted text, or empty string if unsupported/failed.
    """
    if not content:
        return ""

    if len(content) > MAX_ATTACHMENT_SIZE:
        logger.debug(f"Skipping oversized attachment: {filename} ({len(content)} bytes)")
        return ""

    ct = (content_type or '').lower()
    ext = Path(filename).suffix.lower() if filename else ''

    # Try by MIME type first, then by extension
    extractor = _MIME_EXTRACTORS.get(ct) or _EXT_EXTRACTORS.get(ext)
    if not extractor:
        # Fallback: if it looks like text, try plain text
        if ct.startswith('text/') or ext in {'.txt', '.csv', '.log', '.json', '.xml', '.md'}:
            extractor = _extract_plain_text

    if not extractor:
        return ""

    try:
        text = extractor(content)
        if text:
            text = text.strip()
            logger.debug(f"Extracted {len(text)} chars from {filename}")
        return text or ""
    except Exception as e:
        logger.debug(f"Failed to extract text from {filename}: {e}")
        return ""


# --- Extractor implementations ---

def _extract_pdf(content: bytes) -> str:
    """Extract text from PDF."""
    from PyPDF2 import PdfReader
    reader = PdfReader(io.BytesIO(content))
    parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return '\n'.join(parts)


def _extract_docx(content: bytes) -> str:
    """Extract text from .docx (Word)."""
    from docx import Document
    doc = Document(io.BytesIO(content))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Also extract from tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(' | '.join(cells))
    return '\n'.join(parts)


def _extract_xlsx(content: bytes) -> str:
    """Extract text from .xlsx (Excel)."""
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        parts.append(f"[Sheet: {sheet}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(' | '.join(cells))
    wb.close()
    return '\n'.join(parts)


def _extract_pptx(content: bytes) -> str:
    """Extract text from .pptx (PowerPoint)."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            parts.append(f"[Slide {i}]")
            parts.extend(slide_texts)
    return '\n'.join(parts)


def _extract_plain_text(content: bytes) -> str:
    """Extract plain text (txt, csv, log, etc.)."""
    for encoding in ('utf-8', 'latin-1', 'cp1252'):
        try:
            return content.decode(encoding)
        except (UnicodeDecodeError, ValueError):
            continue
    return content.decode('utf-8', errors='replace')


def _extract_html(content: bytes) -> str:
    """Extract text from HTML attachment."""
    import html2text
    text = content.decode('utf-8', errors='replace')
    converter = html2text.HTML2Text()
    converter.ignore_links = True
    converter.ignore_images = True
    converter.body_width = 0
    return converter.handle(text)


# --- Register extractors ---

# PDF
_MIME_EXTRACTORS['application/pdf'] = _extract_pdf
_EXT_EXTRACTORS['.pdf'] = _extract_pdf

# Word
_MIME_EXTRACTORS['application/vnd.openxmlformats-officedocument.wordprocessingml.document'] = _extract_docx
_MIME_EXTRACTORS['application/msword'] = _extract_docx  # .doc — best effort via docx
_EXT_EXTRACTORS['.docx'] = _extract_docx
_EXT_EXTRACTORS['.doc'] = _extract_docx

# Excel
_MIME_EXTRACTORS['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'] = _extract_xlsx
_MIME_EXTRACTORS['application/vnd.ms-excel'] = _extract_xlsx
_EXT_EXTRACTORS['.xlsx'] = _extract_xlsx
_EXT_EXTRACTORS['.xls'] = _extract_xlsx

# PowerPoint
_MIME_EXTRACTORS['application/vnd.openxmlformats-officedocument.presentationml.presentation'] = _extract_pptx
_MIME_EXTRACTORS['application/vnd.ms-powerpoint'] = _extract_pptx
_EXT_EXTRACTORS['.pptx'] = _extract_pptx
_EXT_EXTRACTORS['.ppt'] = _extract_pptx

# Text
for _mime in ('text/plain', 'text/csv'):
    _MIME_EXTRACTORS[_mime] = _extract_plain_text
for _ext in ('.txt', '.csv', '.log', '.json', '.xml', '.md'):
    _EXT_EXTRACTORS[_ext] = _extract_plain_text

# HTML
_MIME_EXTRACTORS['text/html'] = _extract_html
_EXT_EXTRACTORS['.html'] = _extract_html
_EXT_EXTRACTORS['.htm'] = _extract_html
